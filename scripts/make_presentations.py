from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x00, 0x2B, 0x5B)   # header / title background
MID_BLUE    = RGBColor(0x00, 0x5C, 0xA8)   # accent
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)   # subtle background
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY   = RGBColor(0x33, 0x33, 0x33)
MID_GREY    = RGBColor(0x66, 0x66, 0x66)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank = prs.slide_layouts[6]          # truly blank
    return prs.slides.add_slide(blank)


def bg_rect(slide, color):
    """Fill entire slide with a solid colour."""
    sp = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, SLIDE_H
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DARK_GREY,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=16, color=DARK_GREY, bold_first=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size  = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold  = bold_first and first
        first = False
    return txb


def title_slide(prs, title, subtitle, chains):
    slide = blank_slide(prs)
    bg_rect(slide, DARK_BLUE)

    # decorative accent bar
    bar = slide.shapes.add_shape(1, 0, Inches(5.8), SLIDE_W, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = MID_BLUE
    bar.line.fill.background()

    # MTIC label top-left
    add_textbox(slide, "MINISTRY OF TRADE, INDUSTRY AND COOPERATIVES",
                Inches(0.5), Inches(0.3), Inches(10), Inches(0.5),
                font_size=11, color=LIGHT_BLUE)
    add_textbox(slide, "Uganda-UNIDO Programme for Country Partnership (PCP)",
                Inches(0.5), Inches(0.7), Inches(10), Inches(0.4),
                font_size=11, color=LIGHT_BLUE, italic=True)

    # Main title
    add_textbox(slide, title,
                Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0),
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, subtitle,
                Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.6),
                font_size=18, color=LIGHT_BLUE)

    # Chain list
    for i, chain in enumerate(chains):
        add_textbox(slide, f"  {i+1}.  {chain}",
                    Inches(0.5), Inches(4.2) + Inches(i * 0.38),
                    Inches(10), Inches(0.4),
                    font_size=15, color=WHITE)

    # Footer
    add_textbox(slide,
                "Solomon Ariho  |  27 May 2026  |  Draft — for TWG review",
                Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                font_size=11, color=MID_GREY, align=PP_ALIGN.LEFT)


def section_divider(prs, section_number, section_title, subtitle=""):
    slide = blank_slide(prs)
    bg_rect(slide, MID_BLUE)

    bar = slide.shapes.add_shape(1, 0, Inches(3.5), SLIDE_W, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()

    add_textbox(slide, f"Section {section_number}",
                Inches(0.6), Inches(2.2), Inches(11), Inches(0.6),
                font_size=22, color=WHITE, bold=False)
    add_textbox(slide, section_title,
                Inches(0.6), Inches(2.9), Inches(11), Inches(1.2),
                font_size=40, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.6), Inches(4.2), Inches(11), Inches(0.6),
                    font_size=18, color=LIGHT_BLUE, italic=True)


def content_slide(prs, title, bullets, two_col=False, col_split=None):
    slide = blank_slide(prs)
    bg_rect(slide, WHITE)

    # Header bar
    hdr = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = DARK_BLUE
    hdr.line.fill.background()

    # Slide title
    add_textbox(slide, title,
                Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.75),
                font_size=24, bold=True, color=WHITE)

    # Accent line below header
    acc = slide.shapes.add_shape(1, 0, Inches(1.1), SLIDE_W, Inches(0.05))
    acc.fill.solid(); acc.fill.fore_color.rgb = MID_BLUE
    acc.line.fill.background()

    if two_col and col_split:
        mid = col_split
        add_bullet_box(slide, bullets[:mid],
                       Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.8))
        add_bullet_box(slide, bullets[mid:],
                       Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.8))
    else:
        add_bullet_box(slide, bullets,
                       Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.8))

    # Footer rule
    rule = slide.shapes.add_shape(1, 0, Inches(7.25), SLIDE_W, Inches(0.03))
    rule.fill.solid(); rule.fill.fore_color.rgb = LIGHT_BLUE
    rule.line.fill.background()


def table_slide(prs, title, headers, rows, col_widths=None):
    slide = blank_slide(prs)
    bg_rect(slide, WHITE)

    # Header bar
    hdr = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = DARK_BLUE
    hdr.line.fill.background()
    add_textbox(slide, title,
                Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.75),
                font_size=24, bold=True, color=WHITE)

    acc = slide.shapes.add_shape(1, 0, Inches(1.1), SLIDE_W, Inches(0.05))
    acc.fill.solid(); acc.fill.fore_color.rgb = MID_BLUE
    acc.line.fill.background()

    n_cols = len(headers)
    n_rows = len(rows) + 1
    if col_widths is None:
        col_widths = [Inches(13.33 / n_cols)] * n_cols

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.3), Inches(1.3),
        sum(col_widths), Inches(5.9)
    ).table

    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    def set_cell(cell, text, bold=False, bg=None, font_size=13):
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            run = p.runs[0]
            run.font.size  = Pt(font_size)
            run.font.bold  = bold
            run.font.color.rgb = WHITE if bg == DARK_BLUE else DARK_GREY
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    for j, h in enumerate(headers):
        set_cell(tbl.cell(0, j), h, bold=True, bg=DARK_BLUE, font_size=13)

    for i, row in enumerate(rows):
        bg = LIGHT_BLUE if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            set_cell(tbl.cell(i + 1, j), val, bg=bg, font_size=12)


def gantt_slide(prs, title, tasks):
    """
    tasks: list of (task_name, col_idx) where col_idx maps to:
      0=May29  1=May31  2=Jun1  3=Jun2  4=Jun3  5=Jun4
      6=Jun5   7=Jun6   8=Jun7  9=Jun8
    """
    slide = blank_slide(prs)
    bg_rect(slide, WHITE)

    hdr = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = DARK_BLUE
    hdr.line.fill.background()
    add_textbox(slide, title,
                Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.75),
                font_size=24, bold=True, color=WHITE)
    acc = slide.shapes.add_shape(1, 0, Inches(1.1), SLIDE_W, Inches(0.05))
    acc.fill.solid(); acc.fill.fore_color.rgb = MID_BLUE
    acc.line.fill.background()

    DATES    = ['May 29','May 31','Jun 1','Jun 2','Jun 3','Jun 4','Jun 5','Jun 6','Jun 7','Jun 8']
    N_COLS   = len(DATES)
    LEFT     = Inches(2.9)
    TOP      = Inches(1.55)
    RIGHT    = Inches(13.05)
    BOTTOM   = Inches(7.3)
    COL_W    = (RIGHT - LEFT) / N_COLS
    ROW_H    = (BOTTOM - TOP) / len(tasks)

    # Column header labels
    for i, d in enumerate(DATES):
        add_textbox(slide, d,
                    LEFT + i * COL_W, TOP - Inches(0.22),
                    COL_W, Inches(0.2),
                    font_size=8, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Vertical grid lines
    for i in range(N_COLS + 1):
        x = LEFT + i * COL_W
        gl = slide.shapes.add_shape(1, x, TOP, Inches(0.008), BOTTOM - TOP)
        gl.fill.solid(); gl.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        gl.line.fill.background()

    # Rows
    for i, (name, col) in enumerate(tasks):
        y = TOP + i * ROW_H

        # Alternating row background
        if i % 2 == 0:
            rb = slide.shapes.add_shape(1, Inches(0.25), y, RIGHT - Inches(0.25), ROW_H)
            rb.fill.solid(); rb.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xFB)
            rb.line.fill.background()

        # Task label
        add_textbox(slide, name,
                    Inches(0.3), y + Inches(0.02),
                    Inches(2.55), ROW_H - Inches(0.04),
                    font_size=9, color=DARK_GREY)

        # Bar
        pad = ROW_H * 0.18
        bar = slide.shapes.add_shape(
            1,
            LEFT + col * COL_W + COL_W * 0.06,
            y + pad,
            COL_W * 0.88,
            ROW_H - 2 * pad
        )
        bar.fill.solid(); bar.fill.fore_color.rgb = MID_BLUE
        bar.line.fill.background()

    # Footer rule
    rule = slide.shapes.add_shape(1, 0, Inches(7.35), SLIDE_W, Inches(0.03))
    rule.fill.solid(); rule.fill.fore_color.rgb = LIGHT_BLUE
    rule.line.fill.background()


# ════════════════════════════════════════════════════════════════
#  REPORT 1
# ════════════════════════════════════════════════════════════════

def build_report1():
    prs = new_prs()

    # ── Title slide ──────────────────────────────────────────────
    title_slide(prs,
        "Inception Report",
        "Iron & Steel  ·  Copper & Allied Metals  ·  Automotive Value Chains",
        ["Iron & Steel", "Copper & Allied Metals", "Automotives"])

    # ── Section 1: Introduction ──────────────────────────────────
    section_divider(prs, 1, "Introduction")
    content_slide(prs, "Study Overview", [
        "This Inception Report covers three priority value chains designated under NDP IV:",
        "",
        "  1.  Iron & Steel",
        "  2.  Copper & Allied Metals",
        "  3.  Automotives",
        "",
        "Together, these chains represent Uganda's most significant opportunity for mineral-based industrialization —",
        "anchored in the country's iron ore and copper resource endowments and the strategic push toward",
        "domestic vehicle manufacturing through Kiira Motors Corporation.",
        "",
        "This document sets out the study design and workplan, diagnostic framework, data plan,",
        "stakeholder consultation approach, and prioritization methodology.",
    ])

    # ── Section 2: Study Design ──────────────────────────────────
    section_divider(prs, 2, "Study Design & Workplan")

    table_slide(prs, "Report Structure",
        ["Chapter", "Title"],
        [
            ["1", "Executive Summary (written last)"],
            ["2", "Background and Policy Context"],
            ["3", "Methodology and Analytical Framework"],
            ["4", "Iron & Steel Value Chain"],
            ["5", "Copper & Allied Metals Value Chain"],
            ["6", "Automotive Value Chain"],
            ["7", "Cross-Cutting Priority Actions and Policy Recommendations"],
            ["8", "Investment Portfolio and Implementation Roadmap"],
        ],
        col_widths=[Inches(1.2), Inches(10.8)]
    )

    table_slide(prs, "Workplan and Timeline",
        ["Task", "Target Date"],
        [
            ["Inception Report submitted for TWG review", "29 May 2026"],
            ["TWG approval of Inception Report", "31 May 2026"],
            ["Chapter 2: Background and Policy Context", "31 May 2026"],
            ["Chapter 3: Methodology and Analytical Framework", "1 June 2026"],
            ["Chapter 4: Iron & Steel — draft", "2 June 2026"],
            ["Chapter 5: Copper & Allied Metals — draft", "3 June 2026"],
            ["Chapter 6: Automotive — draft", "4 June 2026"],
            ["Chapter 7: Cross-cutting recommendations", "5 June 2026"],
            ["Chapter 8: Investment portfolio and roadmap", "6 June 2026"],
            ["Chapter 1: Executive Summary", "7 June 2026"],
            ["Final report submitted to MTIC", "8 June 2026"],
        ],
        col_widths=[Inches(9.5), Inches(3.0)]
    )

    gantt_slide(prs, "Schedule — Gantt Chart", [
        ("Inception Report draft",    0),
        ("TWG approval",              0),
        ("Ch2 Background & Policy",   1),
        ("Ch3 Methodology",           2),
        ("Ch4 Iron & Steel",          3),
        ("Ch5 Copper & Allied Metals",4),
        ("Ch6 Automotive",            5),
        ("Ch7 Cross-cutting",         6),
        ("Ch8 Investment Roadmap",    7),
        ("Ch1 Executive Summary",     8),
        ("Final submission",          9),
    ])

    # ── Section 3: Diagnostic Framework ─────────────────────────
    section_divider(prs, 3, "Diagnostic Framework", "Six-part template applied across all three chains")

    content_slide(prs, "Analytical Framework — Six Parts", [
        "Part A — Value Chain Map",
        "          Full mapping from primary inputs to all end products. Identifies where Uganda currently participates,",
        "          where value is lost, and where import substitution and export opportunities exist.",
        "",
        "Part B — Current State Assessment",
        "          Production volumes, firms, employment, capacity utilization, export and import values.",
        "",
        "Part C — Binding Constraints Analysis",
        "          Ten dimensions: inputs, technology, energy, water, logistics, waste management,",
        "          finance, skills, standards, and policy/regulatory environment.",
        "",
        "Part D — Market Assessment",
        "          Domestic, regional (EAC/COMESA), and global markets.",
        "",
        "Part E — Prioritized Products",
        "          3–4 products per chain selected via the scoring framework. Explicit deprioritization with rationale.",
        "",
        "Part F — Priority Action Matrix",
        "          Quick wins (0–12 months) · Policy reforms (1–3 years) · Investment interventions (3–5 years)",
    ])

    # ── Section 4: Data Plan ─────────────────────────────────────
    section_divider(prs, 4, "Data Plan")

    table_slide(prs, "Primary Data Sources",
        ["Value Chain", "Primary Sources"],
        [
            ["Iron & Steel", "NPA/UDC Report on the Mapping and Value Chain Analysis (June 2025); KIIs with 40 steel firms, 5 mining players, 9 MDAs"],
            ["Copper & Allied Metals", "KIIs with copper wire/cable manufacturers; UNMC on Kilembe mine; MEMD on copper reserves and licenses"],
            ["Automotives", "KIIs with Kiira Motors Corporation (KMC); bus body fabricators; spare parts manufacturers; MWT on vehicle import data"],
        ],
        col_widths=[Inches(2.8), Inches(10.0)]
    )

    table_slide(prs, "Data Gaps and Mitigation",
        ["Gap", "Chain", "Mitigation"],
        [
            ["Firm-level production data for copper processors", "Copper & Allied Metals", "KIIs with UMA member firms; URA HS code trade data as proxy"],
            ["Copper reserve quantity and grade data", "Copper & Allied Metals", "MEMD consultation; UNMC data on Kilembe"],
            ["Automotive assembly and body fabrication capacity", "Automotives", "KIIs with KMC and private assemblers; URA vehicle import HS codes"],
            ["Regional copper and automotive market sizing", "Both", "UN Comtrade; ITC Trade Map; African Development Bank sector reports"],
        ],
        col_widths=[Inches(4.5), Inches(3.0), Inches(5.5)]
    )

    # ── Section 5: Stakeholder Consultation ─────────────────────
    section_divider(prs, 5, "Stakeholder Consultation Plan")

    table_slide(prs, "Key Informant Interviews (KIIs)",
        ["Stakeholder Category", "Institutions / Firms", "Chain(s)"],
        [
            ["Public sector", "MTIC, UIA, URA, UBOS, UNBS, UEPB, MEMD, UNMC, KMC, MWT", "All"],
            ["Private sector — Iron & Steel", "40 steel firms, SINO Minerals, Tembo Steel, Abyssinia Iron & Steel, Roofings Ltd, Steel & Tube, Uganda Baati", "Iron & Steel"],
            ["Private sector — Copper", "Copper wire/cable manufacturers; scrap copper processors; construction materials importers", "Copper & Allied Metals"],
            ["Private sector — Automotive", "KMC; bus body fabricators; spare parts manufacturers and importers", "Automotives"],
            ["Industry associations", "Uganda Manufacturers Association (UMA)", "All"],
            ["Financial Institutions", "Uganda Development Bank (UDB), Uganda Development Corporation (UDC), East African Development Bank (EADB), commercial banks", "All"],
            ["Research / academia", "Makerere University (engineering, materials science), UIRI", "All"],
        ],
        col_widths=[Inches(3.0), Inches(7.5), Inches(2.5)]
    )

    content_slide(prs, "Validation Workshop", [
        "A structured stakeholder validation workshop will be held before the report is finalized.",
        "",
        "The workshop will:",
        "  ·  Present draft prioritization findings for each of the three value chains",
        "  ·  Test priority product selections and action matrices against stakeholder knowledge",
        "  ·  Identify material gaps or contested findings",
        "  ·  Produce agreed outputs reflected in the final report",
        "",
        "Proposed timing:         Week of 2–6 June 2026",
        "                                  (after chapter drafts are complete, before final submission)",
        "",
        "Proposed participants: TWG members · Private sector representatives · MTIC · UIA · UNBS · UDB",
        "                                  and relevant development partners",
    ])

    # ── Section 6: Prioritization Methodology ───────────────────
    section_divider(prs, 6, "Prioritization Methodology")

    table_slide(prs, "Scoring Framework — Five Criteria",
        ["#", "Criterion", "Description", "Weight"],
        [
            ["1", "Accessible market size", "Size and growth rate of domestic, regional, and export market Uganda can realistically access", "25%"],
            ["2", "Uganda's comparative advantage", "Availability of raw material inputs, existing productive base, geographic proximity to key markets", "25%"],
            ["3", "Feasibility", "Capital requirements, technology readiness, lead time to commercial scale, management capability", "20%"],
            ["4", "Job creation and income potential", "Direct and indirect employment, wage levels, upstream and downstream linkage effects", "15%"],
            ["5", "Import substitution impact", "Value of current imports that domestic production could displace; current account benefit", "15%"],
        ],
        col_widths=[Inches(0.4), Inches(2.6), Inches(8.3), Inches(1.0)]
    )

    table_slide(prs, "Scoring Scale and Interpretation",
        ["Score", "Meaning"],
        [
            ["5", "Excellent — strong evidence of competitive advantage or large accessible market"],
            ["4", "Good — solid case with manageable constraints"],
            ["3", "Moderate — feasible but with significant constraints to address"],
            ["2", "Weak — limited advantage or small market; would require disproportionate investment"],
            ["1", "Poor — Uganda has no realistic competitive position; not recommended"],
        ],
        col_widths=[Inches(1.2), Inches(11.8)]
    )

    content_slide(prs, "Analytical Principle", [
        "Successful countries are known for three or four products.",
        "The economies are built around those few.",
        "Uganda must not try to do everything.",
        "",
        "",
        "",
        "",
        "Composite score = sum of (criterion score × criterion weight) across all five criteria.",
        "",
        "  ·  Above 3.5   →   Priority product",
        "  ·  2.5 – 3.5   →   Conditional priority (subject to enabling conditions)",
        "  ·  Below 2.5   →   Explicitly deprioritized with rationale",
        "",
        "The output of this study is a focused prioritization, not a catalogue.",
        "Every chain analysis must conclude with a clear statement of what Uganda should do",
        "— and what it should NOT do — at this stage.",
    ])

    # ── Section 7: TWG Approval ──────────────────────────────────
    section_divider(prs, 7, "Request for TWG Approval")

    content_slide(prs, "TWG is Requested to:", [
        "1.   Confirm the study scope, diagnostic framework, and assessment templates (Section 3)",
        "",
        "2.   Confirm or revise the prioritization criteria and weights (Section 6)",
        "",
        "3.   Approve the workplan and timeline in Section 2,",
        "      including the 8 June 2026 final submission deadline",
        "",
        "4.   Provide any additional guidance before chapter drafting commences",
        "",
        "",
        "Upon TWG approval, chapter drafting will proceed in accordance with the workplan.",
        "The TWG will be kept informed of progress and any material data gaps as they arise.",
    ])

    prs.save(r"d:\graphics\Jero\MTIC\industrial_diagnostic_study\report\inception-report-1.pptx")
    print("Report 1 saved.")


# ════════════════════════════════════════════════════════════════
#  REPORT 2
# ════════════════════════════════════════════════════════════════

def build_report2():
    prs = new_prs()

    # ── Title slide ──────────────────────────────────────────────
    title_slide(prs,
        "Inception Report",
        "Textiles & Garments  ·  Pharmaceuticals  ·  Petrochemicals & Fertilizers\n"
        "Sugar & Confectionery  ·  Plastics & Packaging  ·  Cement & Building Materials",
        ["Textiles & Garments", "Pharmaceuticals", "Petrochemicals & Fertilizers",
         "Sugar & Confectionery", "Plastics & Packaging", "Cement & Building Materials"])

    # ── Section 1: Introduction ──────────────────────────────────
    section_divider(prs, 1, "Introduction")
    content_slide(prs, "Study Overview", [
        "This Inception Report covers six priority value chains designated under NDP IV:",
        "",
        "  1.  Textiles & Garments              4.  Sugar & Confectionery",
        "  2.  Pharmaceuticals                    5.  Plastics & Packaging",
        "  3.  Petrochemicals & Fertilizers   6.  Cement & Building Materials",
        "",
        "Together, these chains represent Uganda's most significant opportunity for agro-industrial",
        "transformation and manufacturing import substitution — anchored in the country's cotton and",
        "agricultural commodity base, proven capacity in sugar and cement production, and the emerging",
        "potential of domestic pharmaceutical manufacturing and petrochemical development linked to",
        "Uganda's oil and gas resources.",
        "",
        "This document sets out the study design and workplan, diagnostic framework, data plan,",
        "stakeholder consultation approach, and prioritization methodology.",
    ])

    # ── Section 2: Study Design ──────────────────────────────────
    section_divider(prs, 2, "Study Design & Workplan")

    table_slide(prs, "Report Structure",
        ["Chapter", "Title"],
        [
            ["1",  "Executive Summary (written last)"],
            ["2",  "Background and Policy Context"],
            ["3",  "Methodology and Analytical Framework"],
            ["4",  "Textiles & Garments Value Chain"],
            ["5",  "Pharmaceuticals Value Chain"],
            ["6",  "Petrochemicals & Fertilizers Value Chain"],
            ["7",  "Sugar & Confectionery Value Chain"],
            ["8",  "Plastics & Packaging Value Chain"],
            ["9",  "Cement & Building Materials Value Chain"],
            ["10", "Cross-Cutting Priority Actions and Policy Recommendations"],
            ["11", "Investment Portfolio and Implementation Roadmap"],
        ],
        col_widths=[Inches(1.2), Inches(10.8)]
    )

    table_slide(prs, "Workplan and Timeline",
        ["Task", "Target Date"],
        [
            ["Inception Report submitted for TWG review", "29 May 2026"],
            ["TWG approval of Inception Report", "31 May 2026"],
            ["Chapter 2: Background and Policy Context", "31 May 2026"],
            ["Chapter 3: Methodology and Analytical Framework", "1 June 2026"],
            ["Chapter 4: Textiles & Garments — draft", "2 June 2026"],
            ["Chapter 5: Pharmaceuticals — draft", "3 June 2026"],
            ["Chapter 6: Petrochemicals & Fertilizers — draft", "4 June 2026"],
            ["Chapter 7: Sugar & Confectionery — draft", "5 June 2026"],
            ["Chapter 8: Plastics & Packaging — draft", "6 June 2026"],
            ["Chapter 9: Cement & Building Materials — draft", "6 June 2026"],
            ["Chapter 10: Cross-cutting recommendations", "7 June 2026"],
            ["Chapter 11: Investment portfolio and roadmap", "7 June 2026"],
            ["Chapter 1: Executive Summary", "8 June 2026"],
            ["Final report submitted to MTIC", "8 June 2026"],
        ],
        col_widths=[Inches(9.5), Inches(3.0)]
    )

    gantt_slide(prs, "Schedule — Gantt Chart", [
        ("Inception Report draft",       0),
        ("TWG approval",                 0),
        ("Ch2 Background & Policy",      1),
        ("Ch3 Methodology",              2),
        ("Ch4 Textiles & Garments",      3),
        ("Ch5 Pharmaceuticals",          4),
        ("Ch6 Petrochemicals & Fert.",   5),
        ("Ch7 Sugar & Confectionery",    6),
        ("Ch8 Plastics & Packaging",     7),
        ("Ch9 Cement & Build. Mat.",     7),
        ("Ch10 Cross-cutting",           8),
        ("Ch11 Investment Roadmap",      8),
        ("Ch1 Executive Summary",        9),
        ("Final submission",             9),
    ])

    # ── Section 3: Diagnostic Framework ─────────────────────────
    section_divider(prs, 3, "Diagnostic Framework", "Six-part template applied across all six chains")

    content_slide(prs, "Analytical Framework — Six Parts", [
        "Part A — Value Chain Map",
        "          Full mapping from primary inputs to all end products. Identifies where Uganda currently participates,",
        "          where value is lost, and where import substitution and export opportunities exist.",
        "",
        "Part B — Current State Assessment",
        "          Production volumes, firms, employment, capacity utilization, export and import values.",
        "",
        "Part C — Binding Constraints Analysis",
        "          Ten dimensions: inputs, technology, energy, water, logistics, waste management,",
        "          finance, skills, standards, and policy/regulatory environment.",
        "",
        "Part D — Market Assessment",
        "          Domestic, regional (EAC/COMESA), and global markets.",
        "",
        "Part E — Prioritized Products",
        "          3–4 products per chain selected via the scoring framework. Explicit deprioritization with rationale.",
        "",
        "Part F — Priority Action Matrix",
        "          Quick wins (0–12 months) · Policy reforms (1–3 years) · Investment interventions (3–5 years)",
    ])

    # ── Section 4: Data Plan ─────────────────────────────────────
    section_divider(prs, 4, "Data Plan")

    table_slide(prs, "Primary Data Sources",
        ["Value Chain", "Primary Sources"],
        [
            ["Textiles & Garments", "KIIs with textile mills, garment manufacturers, cotton ginneries; CDO data; URA trade data by HS code"],
            ["Pharmaceuticals", "KIIs with pharmaceutical manufacturers; National Drug Authority (NDA) data; URA pharmaceutical import data"],
            ["Petrochemicals & Fertilizers", "KIIs with fertilizer blenders; Ministry of Energy data on oil and gas feedstocks; URA chemical import data"],
            ["Sugar & Confectionery", "KIIs with sugar manufacturers (Kakira, SCOUL, others); Uganda Sugar Board data; URA sugar trade data"],
            ["Plastics & Packaging", "KIIs with plastics manufacturers and recyclers; URA plastic import data; NEMA data on plastic waste"],
            ["Cement & Building Materials", "KIIs with cement plants; Uganda Cement Industry Association data; URA cement trade data; MEMD limestone data"],
        ],
        col_widths=[Inches(3.0), Inches(10.0)]
    )

    table_slide(prs, "Data Gaps and Mitigation",
        ["Gap", "Chain(s)", "Mitigation"],
        [
            ["Firm-level production capacity data", "All six chains", "KIIs with UMA member firms; URA HS code trade data as proxy"],
            ["Cotton lint processing capacity vs. installed capacity", "Textiles & Garments", "CDO consultation; UMA textile members"],
            ["Pharmaceutical manufacturing output by product type", "Pharmaceuticals", "NDA manufacturer register; KIIs with top 10 manufacturers"],
            ["Petrochemical feedstock availability timeline", "Petrochemicals & Fertilizers", "Ministry of Energy and TOTAL/CNOOC/Tullow project data"],
            ["Plastics recycling capacity and waste volumes", "Plastics & Packaging", "NEMA; KIIs with recyclers"],
            ["Limestone reserve quality and location data", "Cement & Building Materials", "MEMD geological survey data"],
        ],
        col_widths=[Inches(4.3), Inches(3.0), Inches(5.7)]
    )

    # ── Section 5: Stakeholder Consultation ─────────────────────
    section_divider(prs, 5, "Stakeholder Consultation Plan")

    table_slide(prs, "Key Informant Interviews (KIIs)",
        ["Stakeholder Category", "Institutions / Firms", "Chain(s)"],
        [
            ["Public sector", "MTIC, UIA, URA, UBOS, UNBS, UEPB, NDA, CDO, Uganda Sugar Board, NEMA, Ministry of Energy", "All"],
            ["Private sector — Textiles", "Nytil, Phenix Logistics, garment exporters, cotton ginneries", "Textiles & Garments"],
            ["Private sector — Pharmaceuticals", "Cipla Quality Chemical, Abacus Pharma, Rene Industries, Joint Medical Store", "Pharmaceuticals"],
            ["Private sector — Petrochemicals", "Fertilizer blenders, petroleum distributors, oil sector players", "Petrochemicals & Fertilizers"],
            ["Private sector — Sugar", "Kakira Sugar Works, SCOUL, Kinyara Sugar, Mayuge Sugar", "Sugar & Confectionery"],
            ["Private sector — Plastics", "Crown Beverages, plastics manufacturers in Namanve, recycling firms", "Plastics & Packaging"],
            ["Private sector — Cement", "Hima Cement, Tororo Cement, National Cement Company", "Cement & Building Materials"],
            ["Industry associations", "Uganda Manufacturers Association (UMA), Uganda Pharmaceutical Manufacturers Association", "All"],
            ["Financial Institutions", "Uganda Development Bank (UDB), Uganda Development Corporation (UDC), East African Development Bank (EADB), commercial banks", "All"],
            ["Research / academia", "Makerere University (pharmacy, chemistry, engineering), UIRI", "All"],
        ],
        col_widths=[Inches(3.0), Inches(7.5), Inches(2.5)]
    )

    content_slide(prs, "Validation Workshop", [
        "A structured stakeholder validation workshop will be held before the report is finalized.",
        "",
        "The workshop will:",
        "  ·  Present draft prioritization findings for each of the six value chains",
        "  ·  Test priority product selections and action matrices against stakeholder knowledge",
        "  ·  Identify material gaps or contested findings",
        "  ·  Produce agreed outputs reflected in the final report",
        "",
        "Proposed timing:         Week of 2–6 June 2026",
        "                                  (after chapter drafts are complete, before final submission)",
        "                                  May be combined with the Report 1 validation workshop",
        "",
        "Proposed participants: TWG members · Private sector representatives from each of the six chains",
        "                                  MTIC · UIA · UNBS · NDA · UDB · relevant development partners",
    ])

    # ── Section 6: Prioritization Methodology ───────────────────
    section_divider(prs, 6, "Prioritization Methodology")

    table_slide(prs, "Scoring Framework — Five Criteria",
        ["#", "Criterion", "Description", "Weight"],
        [
            ["1", "Accessible market size", "Size and growth rate of domestic, regional, and export market Uganda can realistically access", "25%"],
            ["2", "Uganda's comparative advantage", "Availability of raw material inputs, existing productive base, geographic proximity to key markets", "25%"],
            ["3", "Feasibility", "Capital requirements, technology readiness, lead time to commercial scale, management capability", "20%"],
            ["4", "Job creation and income potential", "Direct and indirect employment, wage levels, upstream and downstream linkage effects", "15%"],
            ["5", "Import substitution impact", "Value of current imports that domestic production could displace; current account benefit", "15%"],
        ],
        col_widths=[Inches(0.4), Inches(2.6), Inches(8.3), Inches(1.0)]
    )

    table_slide(prs, "Scoring Scale and Interpretation",
        ["Score", "Meaning"],
        [
            ["5", "Excellent — strong evidence of competitive advantage or large accessible market"],
            ["4", "Good — solid case with manageable constraints"],
            ["3", "Moderate — feasible but with significant constraints to address"],
            ["2", "Weak — limited advantage or small market; would require disproportionate investment"],
            ["1", "Poor — Uganda has no realistic competitive position; not recommended"],
        ],
        col_widths=[Inches(1.2), Inches(11.8)]
    )

    content_slide(prs, "Analytical Principle", [
        "Successful countries are known for three or four products.",
        "The economies are built around those few.",
        "Uganda must not try to do everything.",
        "",
        "",
        "",
        "",
        "Composite score = sum of (criterion score × criterion weight) across all five criteria.",
        "",
        "  ·  Above 3.5   →   Priority product",
        "  ·  2.5 – 3.5   →   Conditional priority (subject to enabling conditions)",
        "  ·  Below 2.5   →   Explicitly deprioritized with rationale",
        "",
        "With six value chains in scope, the discipline of focused prioritization is especially important.",
        "The output must be a clear, defensible set of priorities — not a catalogue.",
    ])

    # ── Section 7: TWG Approval ──────────────────────────────────
    section_divider(prs, 7, "Request for TWG Approval")

    content_slide(prs, "TWG is Requested to:", [
        "1.   Confirm the study scope, diagnostic framework, and assessment templates (Section 3)",
        "",
        "2.   Confirm or revise the prioritization criteria and weights (Section 6)",
        "",
        "3.   Approve the workplan and timeline in Section 2,",
        "      including the 8 June 2026 final submission deadline",
        "",
        "4.   Provide any additional guidance before chapter drafting commences",
        "",
        "",
        "Upon TWG approval, chapter drafting will proceed in accordance with the workplan.",
        "The TWG will be kept informed of progress and any material data gaps as they arise.",
    ])

    prs.save(r"d:\graphics\Jero\MTIC\industrial_diagnostic_study\report\inception-report-2.pptx")
    print("Report 2 saved.")


if __name__ == "__main__":
    build_report1()
    build_report2()
    print("Both presentations created successfully.")
